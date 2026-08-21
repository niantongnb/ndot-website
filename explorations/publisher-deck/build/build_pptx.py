"""Build the ndot publisher deck as .pptx.

Two files:
  editable  - the rendered art plate as the slide background, every word a real
              PowerPoint text box in the right place, in the right face
  flat      - the full render per slide, one picture, nothing editable
"""
import json, os, sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

HERE = os.path.dirname(os.path.abspath(__file__))
SLIDES = os.path.join(HERE, "slides")
OUT = sys.argv[1] if len(sys.argv) > 1 else HERE

PX = 9525                      # 1 CSS px at 96dpi, in EMU
W_PX, H_PX = 1280, 720
FACE = {"N": "Newsreader", "M": "IBM Plex Mono"}
ALIGN = {"l": PP_ALIGN.LEFT, "r": PP_ALIGN.RIGHT, "c": PP_ALIGN.CENTER}

data = json.load(open(os.path.join(HERE, "text.json"), encoding="utf-8"))
STYLES, PAGES = data["styles"], data["pages"]


def new_deck():
    prs = Presentation()
    prs.slide_width = Emu(W_PX * PX)
    prs.slide_height = Emu(H_PX * PX)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def full_bleed(slide, path):
    slide.shapes.add_picture(path, 0, 0, Emu(W_PX * PX), Emu(H_PX * PX))


def set_spacing(run, pts):
    """Character spacing: python-pptx has no API for it, so write the attribute."""
    if pts:
        run.font._rPr.set("spc", str(int(round(pts * 100))))


def add_block(slide, b):
    # a little slack on width and height so a substituted face cannot force a
    # wrap that the browser did not have; TOP anchor keeps the first line put
    slack_w = 6 if b["h"] <= b["lh"] * 1.6 else 2
    box = slide.shapes.add_textbox(
        Emu(int((b["x"] - (slack_w if b["a"] == "r" else 0)) * PX)),
        Emu(int(b["y"] * PX)),
        Emu(int((b["w"] + slack_w) * PX)),
        Emu(int((b["h"] + 8) * PX)),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, para in enumerate(b["p"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ALIGN[b["a"]]
        p.line_spacing = Pt(b["lh"])
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        for run in para:
            face, size, bold, italic, colour, spc = STYLES[run["s"]]
            r = p.add_run()
            r.text = run["t"]
            r.font.name = FACE[face]
            r.font.size = Pt(size)
            r.font.bold = bool(bold)
            r.font.italic = bool(italic)
            r.font.color.rgb = RGBColor.from_string(colour)
            set_spacing(r, spc)
    return box


def build(kind, plate_prefix, with_text):
    prs = new_deck()
    for i in range(len(PAGES)):
        s = blank(prs)
        full_bleed(s, os.path.join(SLIDES, "%s-%02d.png" % (plate_prefix, i + 1)))
        if with_text:
            for b in PAGES[i]:
                add_block(s, b)
    path = os.path.join(OUT, kind)
    prs.save(path)
    print("%-46s %6.1f MB  %d slides" % (os.path.basename(path),
                                         os.path.getsize(path) / 1e6, len(prs.slides.__iter__.__self__._sldIdLst)))
    return path


build("NDot_For_Publishers.pptx", "art", True)
build("NDot_For_Publishers_flat.pptx", "full", False)
